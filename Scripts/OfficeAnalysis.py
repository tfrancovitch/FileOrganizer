#!/usr/bin/env python3
"""
OfficeAnalysis.py
Part of: The File Organizer
Version: 1.3.1

Extracts document properties from Office files in the run's inventory CSV:
title, author, created/modified dates, and a type-appropriate content
count (word count for Word, sheet count for Excel, slide count for
PowerPoint).

Modern Office Open XML files (.docx, .xlsx, .pptx) get full metadata via
python-docx/openpyxl/python-pptx, same as before.

Legacy binary formats (.doc, .xls, .ppt) get METADATA ONLY. Tries olefile
parsing the OLE2 Compound Document structure first (the expected format
for a genuine legacy Office file); if that specifically fails because
the file isn't OLE2 at all, falls back to checking for RTF -- confirmed
via direct diagnosis against real data that this is NOT a rare edge
case: 947 of 958 "OLE2 parse failures" in one real test turned out to be
RTF content saved with a .doc extension, a well-known pattern from
older/cross-platform word processors.

Regardless of which format is actually parsed, no content extraction,
no file conversion, no derivative files created from originals. This is
a deliberate scope decision: this program is built around litigation-
discovery preservation, and creating a converted copy of a document
(even a temporary, carefully-staged one) is out of scope for it. A
separate program handles bulk format conversion for personal
convenience; that is not this one.

Every legacy-file row is explicitly tagged ExtractionMode=MetadataOnly
(vs. Full for modern files) -- both in the CSV and in the report's type
breakdown -- so it's clear at a glance why a .doc file has less data
than a .docx sitting next to it, not something to wonder about later.

Requires:
    pip install python-docx openpyxl python-pptx olefile

Usage:
    python OfficeAnalysis.py --csv DuplicateHashInventory.csv --output OfficeInventory.csv --report OfficeReport.txt
"""

import argparse
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from file_organizer_common import run_analysis, to_long_path

try:
    import docx
except ImportError:
    print("ERROR: python-docx is not installed. Run: pip install python-docx", file=sys.stderr)
    sys.exit(1)

try:
    import openpyxl
except ImportError:
    print("ERROR: openpyxl is not installed. Run: pip install openpyxl", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx is not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    import olefile
except ImportError:
    print("ERROR: olefile is not installed. Run: pip install olefile", file=sys.stderr)
    sys.exit(1)

MODERN_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
LEGACY_EXTENSIONS = {".doc", ".xls", ".ppt"}
EXTENSIONS = MODERN_EXTENSIONS | LEGACY_EXTENSIONS
CHECKPOINT_FIELDS = ["Key", "OfficeType", "ExtractionMode", "Title", "Author", "Created", "Modified", "ContentCount", "Error"]

LEGACY_TYPE_LABELS = {".doc": "Word (legacy)", ".xls": "Excel (legacy)", ".ppt": "PowerPoint (legacy)"}


def analyze_docx(path):
    d = docx.Document(to_long_path(path))
    props = d.core_properties
    word_count = sum(len(p.text.split()) for p in d.paragraphs)
    return {
        "OfficeType": "Word",
        "ExtractionMode": "Full",
        "Title": props.title or "",
        "Author": props.author or "",
        "Created": str(props.created) if props.created else "",
        "Modified": str(props.modified) if props.modified else "",
        "ContentCount": str(word_count),
    }


def analyze_xlsx(path):
    wb = openpyxl.load_workbook(to_long_path(path), read_only=True, data_only=True)
    props = wb.properties
    return {
        "OfficeType": "Excel",
        "ExtractionMode": "Full",
        "Title": props.title or "",
        "Author": props.creator or "",
        "Created": str(props.created) if props.created else "",
        "Modified": str(props.modified) if props.modified else "",
        "ContentCount": str(len(wb.sheetnames)),
    }


def analyze_pptx(path):
    prs = Presentation(to_long_path(path))
    props = prs.core_properties
    return {
        "OfficeType": "PowerPoint",
        "ExtractionMode": "Full",
        "Title": props.title or "",
        "Author": props.author or "",
        "Created": str(props.created) if props.created else "",
        "Modified": str(props.modified) if props.modified else "",
        "ContentCount": str(len(prs.slides)),
    }


def _decode_ole_text(value, codepage):
    """olefile can return metadata text fields (title, author, etc.) as
    raw bytes rather than decoded str -- these are legacy files, not
    always UTF-8 internally. Decode using the file's own codepage
    (e.g. 1252 for Windows-1252/Western European) rather than showing a
    raw bytes repr like "b'Name'". Falls back safely if the codepage is
    missing or the encoding turns out to be wrong."""
    if value is None:
        return ""
    if isinstance(value, bytes):
        try:
            return value.decode(f"cp{codepage}" if codepage else "cp1252")
        except (UnicodeDecodeError, LookupError):
            return value.decode("cp1252", errors="replace")
    return str(value)


RTF_INFO_READ_CAP = 65536  # 64KB -- the \info block appears near the top
                            # of the file; matches the project's existing
                            # 64KB partial-read convention elsewhere.


def _rtf_info_field(rtf_text, field_name):
    """Extract a simple text field like {\\title ...} or {\\author ...}
    from an RTF \\info group. Returns "" if not present -- not every RTF
    file includes one."""
    pattern = r"\{\\" + field_name + r"\s?(.*?)\}"
    match = re.search(pattern, rtf_text, re.DOTALL)
    if not match:
        return ""
    raw = match.group(1)
    raw = re.sub(r"\\[a-z]+-?\d*\s?", " ", raw)  # strip inline RTF control words
    return raw.replace("\\", "").strip()


def _rtf_date_field(rtf_text, field_name):
    """Extract a date field like {\\creatim\\yr2010\\mo5\\dy15\\hr14\\min30}."""
    pattern = r"\\" + field_name + r"\\yr(\d+)\\mo(\d+)\\dy(\d+)(?:\\hr(\d+))?(?:\\min(\d+))?"
    match = re.search(pattern, rtf_text)
    if not match:
        return ""
    yr, mo, dy, hr, mn = match.groups()
    try:
        return f"{int(yr):04d}-{int(mo):02d}-{int(dy):02d} {int(hr or 0):02d}:{int(mn or 0):02d}:00"
    except ValueError:
        return ""


def analyze_rtf(path):
    """Metadata only, via lightweight regex parsing of RTF's optional
    \\info group. RTF is a plain-text markup format, not OLE2 -- a
    genuinely different extraction path, not a variant of the OLE2 one.
    Only reads the first RTF_INFO_READ_CAP bytes, since the \\info block
    appears near the top of the file, not scattered through a
    potentially large document body."""
    long_path = to_long_path(path)
    with open(long_path, "rb") as f:
        raw_bytes = f.read(RTF_INFO_READ_CAP)
    try:
        text = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = raw_bytes.decode("latin-1", errors="replace")

    return {
        "OfficeType": "RTF (legacy)",
        "ExtractionMode": "MetadataOnly",
        "Title": _rtf_info_field(text, "title"),
        "Author": _rtf_info_field(text, "author"),
        "Created": _rtf_date_field(text, "creatim"),
        "Modified": _rtf_date_field(text, "revtim"),
        "ContentCount": "",
    }


def analyze_legacy(path):
    """Metadata only. Tries OLE2 first (the expected format for a
    genuine legacy Office file), and if that specifically fails because
    the file isn't OLE2 at all, checks for RTF's signature before giving
    up -- confirmed via direct diagnosis against real data that this is
    NOT a rare edge case: 947 of 958 "OLE2 parse failures" in one real
    test turned out to be RTF content saved with a .doc extension, a
    well-known pattern from older/cross-platform word processors. See
    module docstring for why this is a deliberate scope boundary
    (metadata only, no conversion) either way."""
    ext = Path(path).suffix.lower()
    long_path = to_long_path(path)

    try:
        ole = olefile.OleFileIO(long_path)
    except olefile.olefile.NotOleFileError:
        # Genuinely not OLE2 -- check the confirmed-common alternative
        # before treating this as an unrecoverable error.
        with open(long_path, "rb") as f:
            header = f.read(8)
        if header.startswith(b"{\\rtf"):
            return analyze_rtf(path)
        # Not OLE2, not RTF either -- re-raise so it's still recorded as
        # a real error (Logs\OfficeInventory.errors.txt), not silently
        # swallowed. Worth running Diagnose-LegacyOfficeFiles.py again if
        # this category turns out to be large.
        raise

    try:
        meta = ole.get_metadata()
        codepage = meta.codepage
        # num_words is populated for legacy Word docs; slides for legacy
        # PowerPoint. Legacy Excel doesn't reliably expose a sheet count
        # via this metadata alone -- left blank rather than guessed.
        content_count = ""
        if ext == ".doc" and meta.num_words is not None:
            content_count = str(meta.num_words)
        elif ext == ".ppt" and meta.slides is not None:
            content_count = str(meta.slides)
        return {
            "OfficeType": LEGACY_TYPE_LABELS[ext],
            "ExtractionMode": "MetadataOnly",
            "Title": _decode_ole_text(meta.title, codepage),
            "Author": _decode_ole_text(meta.author, codepage),
            "Created": str(meta.create_time) if meta.create_time else "",
            "Modified": str(meta.last_saved_time) if meta.last_saved_time else "",
            "ContentCount": content_count,
        }
    finally:
        ole.close()


def analyze_office(path):
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return analyze_docx(path)
    elif ext == ".xlsx":
        return analyze_xlsx(path)
    elif ext == ".pptx":
        return analyze_pptx(path)
    elif ext in LEGACY_EXTENSIONS:
        return analyze_legacy(path)
    raise ValueError(f"Unsupported extension: {ext}")


def report_extra(results):
    by_type = {}
    metadata_only_count = 0
    for r in results:
        t = r.get("OfficeType") or "Unknown"
        by_type[t] = by_type.get(t, 0) + 1
        if r.get("ExtractionMode") == "MetadataOnly":
            metadata_only_count += 1
    lines = ["  By type:"]
    for t, c in sorted(by_type.items()):
        lines.append(f"    {t:<20}: {c}")
    if metadata_only_count:
        lines.append("")
        lines.append(f"  Metadata-only (legacy format, no content extraction): {metadata_only_count}")
        lines.append("  (No conversion or derivative files created -- see module docstring.)")
    return lines


def main():
    parser = argparse.ArgumentParser(description="Extract document properties from Office files (.docx/.xlsx/.pptx, plus metadata-only for legacy .doc/.xls/.ppt).")
    parser.add_argument("--csv", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--report")
    parser.add_argument("--force", action="store_true")
    parser.add_argument("--skip-cloud-only", action="store_true")
    args = parser.parse_args()

    run_analysis(
        csv_path=args.csv, output_path=args.output, report_path=args.report,
        extensions=EXTENSIONS,
        checkpoint_fields=CHECKPOINT_FIELDS, analyze_fn=analyze_office,
        force=args.force, skip_cloud_only=args.skip_cloud_only,
        report_title="OFFICE FILE ANALYSIS REPORT", extra_report_lines_fn=report_extra,
    )


if __name__ == "__main__":
    main()
