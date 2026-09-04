#!/usr/bin/env python3
"""
ContentExtraction.py
Part of: The File Organizer
Version: 1.1.1

Extracts the actual TEXT CONTENT of every content-bearing document in
the database-backed analyzer engine -- PDF, Word (.docx), PowerPoint (.pptx), Excel
(.xlsx), plain text (.txt), and Markdown (.md) -- into individual .txt
files, indexed by DB_ID. This is the last Phase 1 (data-gathering) gap:
earlier scripts captured metadata ABOUT documents (author, page count);
this captures what they actually SAY, which is what a future content-based
organization/search/comparison pass (Phase 2) will need to work with.

Output:
    ExtractedText/<hash>.txt  -- one file per document, full extracted text
    ContentIndex.csv          -- DB_ID, FileName, Path, SourceType,
                                 ExtractedTextFile, CharCount, WordCount, Error

    The supported product runtime stores extracted text by the SHA-256 of
    the extracted UTF-8 text, sharded by hash prefix. Identical extracted
    text therefore reuses one artifact. RunCoordinator uses the in-process analyzer engine and persists results
    directly to SQLite.

Requires:
    pip install pdfplumber python-docx openpyxl python-pptx chardet

"""

import hashlib
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
sys.path.insert(0, str(Path(__file__).parent / "Database"))
from file_organizer_common import to_long_path
import fo_text

try:
    import pdfplumber
except ImportError:
    print("ERROR: pdfplumber is not installed. Run: pip install pdfplumber", file=sys.stderr)
    sys.exit(1)

try:
    import docx
except ImportError:
    print("ERROR: python-docx is not installed. Run: pip install python-docx", file=sys.stderr)
    sys.exit(1)

try:
    import openpyxl
except ImportError:
    print("ERROR: openpyxl is not installed. Run: pip install openpyxl", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx is not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    import chardet
except ImportError:
    print("ERROR: chardet is not installed. Run: pip install chardet", file=sys.stderr)
    sys.exit(1)

EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"}
CHECKPOINT_FIELDS = ["Key", "SourceType", "ExtractedTextFile", "TextSha256",
                     "ReusedExisting", "CharCount", "WordCount", "Error"]


#: Extracted artifacts are sharded two levels deep by the first four
#: hex characters of their content hash: <ab>/<cd>/<full>.txt.
#:
#: B5-E.F008 flagged the flat directory as a real scalability problem
#: -- a large project produced one folder with hundreds of thousands of
#: entries in it, which NTFS handles poorly and Explorer handles worse.
#: Two levels of 256 gives 65,536 buckets, so a million extracted
#: documents average ~15 files per directory.
SHARD_DEPTH = 2


def content_to_relpath(text_sha256):
    r"""
ContentExtraction.py
Part of: The File Organizer B6.1

Per-file content extraction implementation used by the database-backed
in-process analyzer engine.

Supported document text is extracted into a content-addressed, two-level
sharded store keyed by SHA-256 of the extracted UTF-8 text. Identical extracted
text from different source paths reuses one artifact. The analyzer returns the
artifact reference, content hash, encoding/counts and error state; SQLite
persistence is owned by the analyzer runtime, not by this module.

Plain-text decoding is BOM-aware and tries strict UTF-8 before probabilistic
encoding detection. This prevents valid UTF-8 from being silently converted to
a plausible but incorrect legacy encoding.

Requires as applicable: pdfplumber, python-docx, openpyxl, python-pptx, chardet.
"""
    digest = text_sha256
    parts = [digest[i * 2:(i + 1) * 2] for i in range(SHARD_DEPTH)]
    return "/".join(parts + ["%s.txt" % digest])


def path_to_filename(path):
    """B4.5's path-addressed name. RETAINED for `storage_mode` fallback.

    Still reachable when a project sets `extraction.storage_mode` to
    'path_addressed', which exists so an operator upgrading mid-project
    is not forced to re-extract everything at once. New projects get
    content addressing, which is the default in migration 006.
    """
    h = hashlib.sha256(str(path).encode("utf-8")).hexdigest()[:16]
    return f"{h}.txt"


def extract_pdf_text(path):
    parts = []
    with pdfplumber.open(to_long_path(path)) as pdf:
        for page in pdf.pages:
            parts.append(page.extract_text() or "")
    return "PDF", "\n\n".join(parts)


def extract_docx_text(path):
    d = docx.Document(to_long_path(path))
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "Word", "\n".join(parts)


def extract_pptx_text(path):
    prs = Presentation(to_long_path(path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts = [
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        ]
        parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))
    return "PowerPoint", "\n\n".join(parts)


def extract_xlsx_text(path):
    wb = openpyxl.load_workbook(to_long_path(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join("" if v is None else str(v) for v in row)
            if line.strip():
                parts.append(line)
    return "Excel", "\n".join(parts)


def extract_plain_text(path):
    with open(to_long_path(path), "rb") as f:
        raw = f.read()
    text, _encoding = fo_text.decode_bytes(raw)
    return "PlainText", text


def make_analyze_fn(extract_folder, content_addressed=True):
    def analyze_content(path):
        ext = Path(path).suffix.lower()
        if ext == ".pdf":
            source_type, text = extract_pdf_text(path)
        elif ext == ".docx":
            source_type, text = extract_docx_text(path)
        elif ext == ".pptx":
            source_type, text = extract_pptx_text(path)
        elif ext == ".xlsx":
            source_type, text = extract_xlsx_text(path)
        elif ext in (".txt", ".md"):
            source_type, text = extract_plain_text(path)
        else:
            raise ValueError(f"Unsupported extension: {ext}")

        # Content addressing: the artifact is named for what is IN it.
        text_sha = hashlib.sha256(text.encode("utf-8")).hexdigest()
        if content_addressed:
            relpath = content_to_relpath(text_sha)
        else:
            relpath = path_to_filename(path)

        target = extract_folder / relpath
        target.parent.mkdir(parents=True, exist_ok=True)

        # A hit here means another document already produced byte-for-byte
        # this text. Writing it again would produce an identical file, so
        # the write is skipped and the reuse is REPORTED rather than
        # hidden -- a consumer counting artifacts should be able to tell
        # deduplication from extraction failure.
        reused = target.exists()
        if not reused:
            # Written to a temporary neighbour and renamed, so a crash
            # mid-write cannot leave a truncated artifact sitting at a
            # content-addressed name that later runs will trust (B5-G).
            staging = target.with_suffix(".txt.partial")
            with open(staging, "w", encoding="utf-8") as f:
                f.write(text)
            os.replace(staging, target)

        # B6: counted, not materialised. len(text.split()) built a list
        # of every token to produce one integer -- B5-E.F009 measured an
        # 8 MB document costing ~110 MB of heap that way.
        return {
            "SourceType": source_type,
            "ExtractedTextFile": relpath,
            "TextSha256": text_sha,
            "ReusedExisting": "True" if reused else "False",
            "CharCount": str(len(text)),
            "WordCount": str(fo_text.count_words(text)),
        }
    return analyze_content


def report_extra(results):
    total_chars = sum(int(r["CharCount"]) for r in results if r.get("CharCount") and not r["Error"])
    total_words = sum(int(r["WordCount"]) for r in results if r.get("WordCount") and not r["Error"])
    empty_count = sum(1 for r in results if r.get("CharCount") == "0")

    by_type = {}
    for r in results:
        t = r.get("SourceType") or "Unknown"
        by_type[t] = by_type.get(t, 0) + 1

    lines = [
        f"  Total characters extracted   : {total_chars:,}",
        f"  Total words extracted        : {total_words:,}",
        f"  No extractable text (likely scanned/empty): {empty_count}",
        "  By source type:",
    ]
    for t, c in sorted(by_type.items(), key=lambda x: -x[1]):
        lines.append(f"    {t:<12}: {c}")
    return lines
